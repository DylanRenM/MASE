"""
Analyze PPTX layout details and compare with HTML structure.
"""
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

pptx_path = '/Users/dylanren/Documents/trae_projects/日常办公/projects/unified-dev-framework-training.pptx'
html_path = '/Users/dylanren/Documents/trae_projects/merge-books/.frontend-slides/unified-dev-framework-training.html'

prs = Presentation(pptx_path)

print("=" * 80)
print("PPTX Analysis Report")
print("=" * 80)

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    print(f"Layout: {slide.slide_layout.name}")
    
    shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text[:50] if shape.text else "(empty)"
            shapes.append({
                'name': shape.name,
                'type': 'textbox',
                'x': shape.left,
                'y': shape.top,
                'width': shape.width,
                'height': shape.height,
                'text': text
            })
        else:
            shapes.append({
                'name': shape.name,
                'type': 'shape',
                'x': shape.left,
                'y': shape.top,
                'width': shape.width,
                'height': shape.height
            })
    
    shapes.sort(key=lambda s: (s['y'], s['x']))
    
    for s in shapes:
        print(f"  {s['type']:10s} x={s['x']/914400:6.2f} y={s['y']/914400:6.2f} w={s['width']/914400:6.2f} h={s['height']/914400:6.2f} | {s['text']}")

print("\n" + "=" * 80)
print("HTML Structure Analysis")
print("=" * 80)

tree = etree.parse(html_path)
root = tree.getroot()

slides = root.xpath('//div[@class="deck-stage"]/section')
print(f"Found {len(slides)} slides in HTML")

for i, slide in enumerate(slides[:3]):
    print(f"\n--- HTML Slide {i+1} ---")
    for child in slide:
        tag = child.tag
        classes = child.get('class', '')
        style = child.get('style', '')
        text = child.text.strip() if child.text and child.text.strip() else ''
        if not text and len(child):
            text = f"{len(child)} children"
        print(f"  {tag:10s} class={classes[:30]:30s} style={style[:50]:50s} | {text[:40]}")
