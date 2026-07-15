"""
Convert HTML slide deck to PPTX using screenshots.
"""
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu
import os

HTML_WIDTH = 1920
HTML_HEIGHT = 1080
PPTX_WIDTH = 1280
PPTX_HEIGHT = 720

def generate_screenshots(html_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={'width': HTML_WIDTH, 'height': HTML_HEIGHT})
        
        page.goto(f'file://{html_path}')
        page.wait_for_timeout(3000)
        
        slide_elements = page.query_selector_all('section.slide')
        
        for i, slide_el in enumerate(slide_elements):
            page.evaluate("""
                document.querySelectorAll('section.slide').forEach(el => {
                    el.style.visibility = 'hidden';
                    el.style.opacity = '0';
                });
            """)
            
            slide_el.evaluate('el => { el.style.visibility = "visible"; el.style.opacity = "1"; }')
            page.wait_for_timeout(500)
            
            screenshot_path = os.path.join(output_dir, f'slide_{i+1:03d}.png')
            page.screenshot(path=screenshot_path, full_page=False)
            
            print(f'Saved screenshot: {screenshot_path}')
        
        browser.close()

def create_pptx(screenshots_dir, output_pptx):
    prs = Presentation()
    
    prs.slide_width = Emu(PPTX_WIDTH * 9525)
    prs.slide_height = Emu(PPTX_HEIGHT * 9525)
    
    screenshot_files = sorted([f for f in os.listdir(screenshots_dir) if f.endswith('.png')])
    
    for i, screenshot_file in enumerate(screenshot_files):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        screenshot_path = os.path.join(screenshots_dir, screenshot_file)
        
        slide.shapes.add_picture(
            screenshot_path,
            Emu(0),
            Emu(0),
            width=Emu(PPTX_WIDTH * 9525),
            height=Emu(PPTX_HEIGHT * 9525)
        )
        
        print(f'Slide {i+1}: Added screenshot')
    
    prs.save(output_pptx)
    print(f'Saved to: {output_pptx}')

def main():
    html_path = '/Users/dylanren/Documents/trae_projects/merge-books/.frontend-slides/unified-dev-framework-training.html'
    screenshots_dir = '/Users/dylanren/Documents/trae_projects/日常办公/projects/screenshots'
    output_pptx = '/Users/dylanren/Documents/trae_projects/日常办公/projects/unified-dev-framework-training.pptx'
    
    print('Generating screenshots...')
    generate_screenshots(html_path, screenshots_dir)
    
    print('Creating PPTX...')
    create_pptx(screenshots_dir, output_pptx)
    
    print('Done!')

if __name__ == '__main__':
    main()
