#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""验证生成的PPT"""

from pptx import Presentation
from pptx.util import Emu

pptx_path = '/Users/dylanren/Documents/trae_projects/日常办公/MASE框架培训讲义.pptx'

prs = Presentation(pptx_path)

print(f"幻灯片尺寸: {prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} 英寸")
print(f"总幻灯片数: {len(prs.slides)}\n")
print("=" * 80)

for i, slide in enumerate(prs.slides, 1):
    print(f"\n第 {i} 张幻灯片:")
    print("-" * 60)
    
    # 收集所有文本
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in para.runs)
                if text.strip():
                    texts.append(text.strip())
    
    # 打印前几条文本
    for j, text in enumerate(texts[:6]):
        print(f"  {text[:80]}")
    if len(texts) > 6:
        print(f"  ... (共 {len(texts)} 段文本)")

print("\n" + "=" * 80)
print("\n验证完成！")
