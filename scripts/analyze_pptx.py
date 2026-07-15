"""
Analyze the generated PPTX to check layout issues.
Reads all shapes and their positions/sizes to identify problems.
"""
from pptx import Presentation

pptx_path = '/Users/dylanren/Documents/trae_projects/日常办公/projects/unified-dev-framework-training.pptx'
prs = Presentation(pptx_path)

print(f"Total slides: {len(prs.slides)}")
print(f"Slide size: {prs.slide_width / 9525:.0f} x {prs.slide_height / 9525:.0f}")
print("=" * 80)

for idx, slide in enumerate(prs.slides):
    print(f"\nSlide {idx+1}:")
    shapes_info = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text[:50] if shape.text_frame.text else ""
            shapes_info.append({
                'name': shape.name,
                'x': shape.left / 9525,
                'y': shape.top / 9525,
                'w': shape.width / 9525,
                'h': shape.height / 9525,
                'text': text,
            })
    
    shapes_info.sort(key=lambda s: (s['y'], s['x']))
    
    for s in shapes_info:
        print(f"  [{s['y']:.1f}px] [{s['x']:.1f}px] [{s['w']:.1f}x{s['h']:.1f}] {s['text'][:60]}")
    
    if idx < 5:
        print("-" * 60)
